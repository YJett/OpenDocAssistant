import docx
import matplotlib.pyplot as plt
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.oxml import parse_xml
from docx.oxml.ns import nsdecls
from src import api_doc,word_reader
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx import Document
from docx.shared import Inches, Pt, RGBColor
import os
import matplotlib
from pptx import Presentation
import logging


# 设置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# 全局变量定义
doc = None
current_paragraph = None 
current_table = None
current_picture = None
current_heading = None
chart = None
# 图表相关的全局常量
CONTENT_WIDTH = Inches(6)  # 内容区域宽度
CONTENT_HEIGHT = Inches(4) # 内容区域高度
CONTENT_LEFT = Inches(1)   # 内容区域左边距
CONTENT_TOP = Inches(1)    # 内容区域上边距

# 图片相关的全局常量
PIC_PATH = "test/pics"  # 图片路径
PIC_LEFT = CONTENT_LEFT  # 图片默认左边距
PIC_TOP = CONTENT_TOP   # 图片默认上边距

def set_word(docx_path=None):
    """设置当前操作的Word文档
    
    Args:
        docx_path: Word文档路径，如果为None则创建新文档
    Returns:
        Document对象
    """
    global doc, current_paragraph, current_table, current_picture, current_heading
    
    try:
        if docx_path and os.path.exists(docx_path):
            doc = Document(docx_path)
            logger.info(f"Opened document: {docx_path}")
        else:
            doc = Document()
            logger.info("Created new document")
            
        # 重置所有当前操作对象
        current_paragraph = None
        current_table = None
        current_picture = None
        current_heading = None
        
        return doc
        
    except Exception as e:
        logger.error(f"Error in set_word: {e}")
        raise

def get_word():
    """获取当前Word文档对象"""
    global doc
    return doc


def create_docx():
    """创建新的Word文档
    
    类似于PPT中的create_slide()函数，但Word不需要创建幻灯片。
    初始化文档和全局状态变量。
    """
    global doc, current_paragraph, current_table, current_picture, current_heading
    try:
        doc = Document()
        logger.info("Created new document")
        
        # 重置所有当前操作对象
        current_paragraph = None
        current_table = None
        current_picture = None
        current_heading = None
        
        # 设置文档阅读器
        word_reader.set_document(doc)
        
        return doc
    except Exception as e:
        logger.error(f"Error creating document: {e}")
        return None

def save_word(docx_path):
    """保存Word文档到指定路径
    
    Args:
        docx_path: 保存路径
    """
    global doc
    try:
        # 确保目录存在
        os.makedirs(os.path.dirname(docx_path), exist_ok=True)
        doc.save(docx_path)
        logger.info(f"Document saved to: {docx_path}")
        return True
    except Exception as e:
        logger.error(f"Error saving document: {e}")
        return False

def save_state():
    """保存当前状态"""
    global doc, current_paragraph, current_table, current_picture, current_heading, chart
    return {
        'doc': doc,
        'current_paragraph': current_paragraph,
        'current_table': current_table,
        'current_picture': current_picture,
        'current_heading': current_heading,
        'chart': chart
    }
    return state

def load_state(state):
    """加载保存的状态"""
    global doc, current_paragraph, current_table, current_picture, current_heading, chart
    doc = state['doc']
    current_paragraph = state['current_paragraph']
    current_table = state['current_table']
    current_picture = state['current_picture']
    current_heading = state['current_heading']
    chart = state['chart']


# 2. 修改 docx 文件的名称（不改变路径）
def rename_docx_name(original_path, new_name):
    try:
        if not os.path.exists(original_path):
            print(f"错误: 文件 {original_path} 不存在！")
            return False
        directory = os.path.dirname(original_path)
        new_path = os.path.join(directory, new_name)
        os.rename(original_path, new_path)
        print(f"文件成功重命名为: {new_path}")
        return True
    except Exception as e:
        print(f"发生错误: {e}")
        return False


# 3. 修改 docx 文件的路径（不改变文件名）
def move_docx_to_new_path(original_path, new_directory):
    try:
        if not os.path.exists(original_path):
            print(f"错误: 文件 {original_path} 不存在！")
            return False
        file_name = os.path.basename(original_path)
        new_path = os.path.join(new_directory, file_name)
        os.rename(original_path, new_path)
        print(f"文件成功移动到新路: {new_path}")
        return True
    except Exception as e:
        print(f"发生错误: {e}")
        return False


# 4. 创建段落并设置格式
def add_paragraph(text, style=None):
    """添加段落
    
    Args:
        text: 段落文本
        style: 段落样式(可选)
    
    Returns:
        新添加的段落对象
    """
    global doc, current_paragraph
    try:
        current_paragraph = doc.add_paragraph(text, style=style)
        return current_paragraph
    except Exception as e:
        logger.error(f"Error adding paragraph: {e}")
        return None


# 6. 根据标题查找段落
def find_paragraphs_by_heading_and_content(heading, content):
    global doc
    paragraphs = []
    found_heading = False
    for para in doc.paragraphs:
        if para.style.name.startswith('Heading') and para.text == heading:
            found_heading = True
            continue
        if found_heading and para.style.name.startswith('Normal') and para.text.strip() == content:
            paragraphs.append(para.text)
        if found_heading and para.style.name.startswith('Heading'):
            break
    if paragraphs:
        print(f"找到了标题 '{heading}' 下匹配内容 '{content}' 的段落: {paragraphs}")
    else:
        print(f"没有找到标题 '{heading}' 下匹配内容 '{content}' 的段落。")
    return paragraphs


def delete_header(section_index=None):
    """
    删除页眉
    如果指定section_index,则只删除该节的页眉
    否则删除所有节的页眉
    """
    global doc
    if section_index is not None:
        if 0 <= section_index < len(doc.sections):
            section = doc.sections[section_index]
            # 断开与前一节的链接
            section.header.is_linked_to_previous = False
            section.header.paragraphs[0].text = ""
    else:
        for section in doc.sections:
            # 断开与前一节的链接
            section.header.is_linked_to_previous = False
            section.header.paragraphs[0].text = ""

def add_header(text, section_index=None):
    """
    添加页眉
    如果指定section_index,则只为该节添加页眉
    否则为所有节添加页眉
    """
    global doc
    if section_index is not None:
        if 0 <= section_index < len(doc.sections):
            section = doc.sections[section_index]
            # 断开与前一节的链接
            section.header.is_linked_to_previous = False
            section.header.paragraphs[0].text = text
    else:
        for i, section in enumerate(doc.sections):
            # 断开与前一节的链接
            section.header.is_linked_to_previous = False
            section.header.paragraphs[0].text = text

def add_footer(text, section_index=None):
    """
    添加页脚
    如果指定section_index,则只为该节添加页脚
    否则为所有节添加页脚
    """
    global doc
    if section_index is not None:
        if 0 <= section_index < len(doc.sections):
            section = doc.sections[section_index]
            # 断开与前一节的链接
            section.footer.is_linked_to_previous = False
            section.footer.paragraphs[0].text = text
    else:
        for section in doc.sections:
            # 断开与前一节的链接
            section.footer.is_linked_to_previous = False
            section.footer.paragraphs[0].text = text

def delete_footer(section_index=None):
    """
    删除页脚
    如果指定section_index,则只删除该节的页脚
    否则删除所有节的页脚
    """
    global doc
    if section_index is not None:
        if 0 <= section_index < len(doc.sections):
            section = doc.sections[section_index]
            # 断开与前一节的链接
            section.footer.is_linked_to_previous = False
            section.footer.paragraphs[0].text = ""
    else:
        for section in doc.sections:
            # 断开与前一节的链接
            section.footer.is_linked_to_previous = False
            section.footer.paragraphs[0].text = ""


def add_line_break(paragraph_index, position=None):
    """
    在指定段落中添加换行符

    参数:
    - paragraph_index: 段落索引 
    - position: 在段落中插入的位置(如果不指定,则在段落末尾添加)
    """
    global doc
    if 0 <= paragraph_index < len(doc.paragraphs):
        paragraph = doc.paragraphs[paragraph_index]

        # 创建新的 run 并移动到指定位置
        new_run = paragraph.add_run()
        if position is not None and 0 <= position < len(paragraph.runs):
            paragraph._p.insert(position, new_run._r)
            paragraph.runs.insert(position, new_run)
            paragraph.runs.pop()

        new_run.add_break(WD_BREAK.LINE)


def delete_line_break(paragraph_index, break_index):
    """
    删除指定段落中的换行符

    参数:
    - paragraph_index: 段落索引
    - break_index: 换行符在段落中的索引
    """
    global doc
    if 0 <= paragraph_index < len(doc.paragraphs):
        paragraph = doc.paragraphs[paragraph_index]
        breaks = paragraph._element.findall('.//w:br', namespaces={
            'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'})
        if 0 <= break_index < len(breaks):
            line_break = breaks[break_index]
            parent = line_break.getparent()
            if parent is not None:
                parent.remove(line_break)
                # 如果 run 为空，则删除该 run
                if not parent.getchildren():
                    parent.getparent().remove(parent)


# def add_style(doc_path, style_name, font_name=None, font_size=None, bold=None, italic=None, color=None):
#     """
#     添加新样式

#     参数:
#     - style_name: 样式名称
#     - font_name: 字名称
#     - font_size: 字体大小(磅)
#     - bold: 是否加粗
#     - italic: 是否斜体
#     - color: RGB颜色元组,如(255, 0, 0)

#     返回: bool 是否添加成功
#     """
#     doc = Document(doc_path)

#     # 检查样式是否已存在
#     if style_name in doc.styles:
#         return False

#     try:
#         style = doc.styles.add_style(style_name, WD_STYLE_TYPE.PARAGRAPH)
#         if font_name:
#             style.font.name = font_name
#         if font_size:
#             style.font.size = Pt(font_size)
#         if bold is not None:
#             style.font.bold = bold
#         if italic is not None:
#             style.font.italic = italic
#         if color:
#             style.font.color.rgb = RGBColor(*color)
#         doc.save(doc_path)
#         return True
#     except Exception:
#         return False


# def delete_style(doc_path, style_name):
#     """
#     删除指定样式

#     参数:
#     - style_name: 要删除的样式名称

#     返回: bool 是否删除成功
#     """
#     doc = Document(doc_path)
#     try:
#         # 将使用样式的段落重置为默认样式
#         for paragraph in doc.paragraphs:
#             if paragraph.style.name == style_name:
#                 paragraph.style = doc.styles['Normal']

#         # 删除样式
#         doc.styles._element.remove(doc.styles[style_name]._element)
#         doc.save(doc_path)
#         return True
#     except Exception:
#         return False


def add_page_numbers():
    """添加页码（为了测试删除页码功能）"""
    global doc
    sections = doc.sections
    for section in sections:
        footer = section.footer
        paragraph = footer.paragraphs[0]
        run = paragraph.add_run()

        # 添加页码字段
        fldChar1 = OxmlElement('w:fldChar')
        fldChar1.set(qn('w:fldCharType'), 'begin')
        run._r.append(fldChar1)

        instrText = OxmlElement('w:instrText')
        instrText.text = "PAGE"
        run._r.append(instrText)

        fldChar2 = OxmlElement('w:fldChar')
        fldChar2.set(qn('w:fldCharType'), 'end')
        run._r.append(fldChar2)

def delete_page_numbers():
    """
    删除页码
    """
    global doc
    for section in doc.sections:
        # 检查页眉
        for paragraph in section.header.paragraphs:
            if 'PAGE' in paragraph._element.xml:
                paragraph._element.getparent().remove(paragraph._element)

        # 检查页脚
        for paragraph in section.footer.paragraphs:
            if 'PAGE' in paragraph._element.xml:
                paragraph._element.getparent().remove(paragraph._element)

def add_table_of_contents(levels=3):
    """
    添加目录

    参数:
    - levels: 目录级别(1-9)
    """
    global doc
    paragraph = doc.add_paragraph()
    run = paragraph.add_run()

    fldChar = OxmlElement('w:fldChar')
    fldChar.set(qn('w:fldCharType'), 'begin')
    run._element.append(fldChar)

    instrText = OxmlElement('w:instrText')
    instrText.text = f'TOC \\o "1-{levels}" \\h \\z \\u'
    run._element.append(instrText)

    fldChar = OxmlElement('w:fldChar')
    fldChar.set(qn('w:fldCharType'), 'end')
    run._element.append(fldChar)


def delete_table_of_contents():
    """
    删除目录
    """
    global doc
    for paragraph in doc.paragraphs:
        if 'TOC' in paragraph._element.xml:
            paragraph._element.getparent().remove(paragraph._element)


def add_watermark(text):
    """添加文字水印"""
    global doc
    for section in doc.sections:
        header = section.header
        paragraph = header.paragraphs[0]
        run = paragraph.add_run()

        # 创建水印
        hdr = parse_xml(f'''
            <w:pict xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" xmlns:v="urn:schemas-microsoft-com:vml">
                <v:shape id="_x0000_i1025" type="#_x0000_t136" style="position:absolute;margin-left:0;margin-top:0;width:468pt;height:351pt;rotation:315;z-index:-251654144">
                    <v:textpath style="font-family:'Calibri';font-size:1pt" string="{text}"/>
                </v:shape>
            </w:pict> 
            ''')

        run._r.append(hdr)

def delete_watermark():
    """
    删除水印
    """
    global doc
    for section in doc.sections:
        header = section.header
        for paragraph in header.paragraphs:
            for run in paragraph.runs:
                if 'pict' in run._element.xml:
                    run._element.getparent().remove(run._element)


def delete_paragraph(paragraph_index):
    """
    删除指定索引的段落
    """
    global doc
    if 0 <= paragraph_index < len(doc.paragraphs):
        p = doc.paragraphs[paragraph_index]._element
        p.getparent().remove(p)


# # 示例使用：
# content = "这是一个新的文档，包一些示例内容。"
# save_path = "new_document_with_headings.docx"
# create_and_save_docx(content, save_path)
#
# doc = Document()
# add_heading(doc, "一级标题", level=1, font_name='Calibri', font_size=Pt(16), color=RGBColor(0, 0, 255), spacing=Pt(1.5))
# add_heading(doc, "二级标题", level=2, font_name='Calibri', font_size=Pt(14), color=RGBColor(255, 0, 0), spacing=Pt(1.2))
# add_paragraph(doc, "这是一个普通的段落。", font_name='Times New Roman', font_size=Pt(11), color=RGBColor(255, 255, 0), spacing=Pt(1.1))
#
# add_heading(doc, "三级标题", level=3, font_name='Calibri', font_size=Pt(12), color=RGBColor(0, 255, 0), spacing=Pt(1.0))
# add_paragraph(doc, "这是一个普通的段落。", font_name='Times New Roman', font_size=Pt(11), color=RGBColor(255, 255, 0), spacing=Pt(1.1))
# add_paragraph(doc, "这是另一个段落，紧接在标题后面。", font_name='Arial', font_size=Pt(10), color=RGBColor(128, 0, 128), spacing=Pt(1.0))
# doc.save("new_document_with_headings.docx")
# content = "这是一普通的段落。"
# found_paragraphs = find_paragraphs_by_heading_and_content(doc, "二级标题",content)


# text


color2hex = {
    "blue": "0000FF", "light blue": "ADD8E6", "dark blue": "00008B",
    "green": "008000", "light green": "90EE90", "dark green": "006400",
    "yellow": "FFFF00", "light yellow": "FFFFE0", "dark yellow": "BDB76B",
    "orange": "FFA500", "light orange": "FFDAB9", "dark orange": "FF8C00",
    "red": "FF0000", "light red": "FFC0CB", "dark red": "8B0000",
    "black": "000000", "white": "FFFFFF", "purple": "800080", "pink": "FFC0CB",
}


def add_text(paragraph_index, text, position=None, bold=False, italic=False, underline=False, color=None,
             size=None):
    global doc
    if 0 <= paragraph_index < len(doc.paragraphs):
        paragraph = doc.paragraphs[paragraph_index]
        if position is not None:
            # 分割现有文本并创建新的运行块
            old_text = paragraph.text
            paragraph.clear()
            if position > 0:
                run = paragraph.add_run(old_text[:position])
            run = paragraph.add_run(text)
            if position < len(old_text):
                paragraph.add_run(old_text[position:])
        else:
            run = paragraph.add_run(text)

        run.bold = bold
        run.italic = italic
        run.underline = underline
        if color:
            run.font.color.rgb = RGBColor(*color)
        if size:
            run.font.size = Pt(size)

def delete_text(paragraph_index, start, end):
    global doc
    if 0 <= paragraph_index < len(doc.paragraphs):
        paragraph = doc.paragraphs[paragraph_index]
        runs = paragraph.runs
        text = paragraph.text

        # 保存现有运行块的格式信息
        formats = []
        current_pos = 0
        for run in runs:
            formats.append({
                'text': run.text,
                'start': current_pos,
                'end': current_pos + len(run.text),
                'bold': run.bold,
                'italic': run.italic,
                'underline': run.underline,
                'color': run.font.color.rgb,
                'size': run.font.size
            })
            current_pos += len(run.text)

        # 创建新文本
        new_text = text[:start] + text[end:]

        # 重新应用格式
        paragraph.clear()
        current_pos = 0
        for fmt in formats:
            if current_pos < len(new_text):
                # 计算新的文本范围
                new_start = max(0, fmt['start'] - (end - start))
                new_end = max(0, fmt['end'] - (end - start))
                if new_start < new_end:
                    run = paragraph.add_run(new_text[new_start:new_end])
                    run.bold = fmt['bold']
                    run.italic = fmt['italic']
                    run.underline = fmt['underline']
                    if fmt['color']:
                        run.font.color.rgb = fmt['color']
                    if fmt['size']:
                        run.font.size = fmt['size']
                    current_pos = new_end

from docx.oxml.shared import OxmlElement, qn


def add_hyperlink(paragraph_index, text, url, position=None):
    """
    在指定段落中添加超链接

    参数:
    - paragraph_index: 段落索引
    - text: 链接显示的文本
    - url: 链接地址
    - position: 在段落中插入的位置(如果不指定,则在段落末尾添加)
    """
    global doc
    if 0 <= paragraph_index < len(doc.paragraphs):
        paragraph = doc.paragraphs[paragraph_index]

        # 创建新的运行块
        run = paragraph.add_run()
        run.text = text
        run.font.underline = True
        run.font.color.rgb = RGBColor(0, 0, 255)  # 蓝色

        # 添加 rStyle
        rStyle = OxmlElement('w:rStyle')
        rStyle.set(qn('w:val'), 'Hyperlink')
        run._element.rPr.append(rStyle)

        # 添加超链接关系
        r_id = doc.part.relate_to(url, docx.opc.constants.RELATIONSHIP_TYPE.HYPERLINK, is_external=True)
        hyperlink = OxmlElement('w:hyperlink')
        hyperlink.set(qn('r:id'), r_id)
        hyperlink.append(run._element)

        # 处理插入位置
        if position is not None and 0 <= position <= len(paragraph._p):
            paragraph._p.insert(position, hyperlink)
        else:
            paragraph._p.append(hyperlink)

def delete_hyperlink(paragraph_index, link_index):
    """
    删除指定段落中的超链接

    参数:
    - paragraph_index: 段落索引
    - link_index: 超链接在段落中的索引
    """
    global doc
    if 0 <= paragraph_index < len(doc.paragraphs):
        paragraph = doc.paragraphs[paragraph_index]
        hyperlinks = paragraph._p.xpath('.//w:hyperlink')
        if 0 <= link_index < len(hyperlinks):
            hyperlink = hyperlinks[link_index]
            for run in hyperlink.xpath('.//w:r'):
                hyperlink.addnext(run)
            hyperlink.getparent().remove(hyperlink)

def add_heading(text, level):
    """添加标题
    
    Args:
        text: 标题文本
        level: 标题级别(1-9)
    Returns:
        新添加的标题段落对象
    """
    global doc, current_paragraph
    try:
        current_paragraph = doc.add_heading(text, level)
        return current_paragraph
    except Exception as e:
        logger.error(f"Error adding heading: {e}")
        return None


def delete_heading(level, occurrence):
    """
    删除指定级别的第n个标题

    参数:
    - level: 标题级别
    - occurrence: 第几个该级别的标题(从1开始)
    """
    global doc
    count = 0
    for i, paragraph in enumerate(doc.paragraphs):
        if paragraph.style.name.startswith(f'Heading {level}'):
            count += 1
            if count == occurrence:
                p = doc.paragraphs[i]._element
                p.getparent().remove(p)
                break

# 设置字体大小
def set_font_size(size):
    """设置字体大小"""
    global current_paragraph
    for run in current_paragraph.runs:
        run.font.size = Pt(size)


# 设置字体颜色
def set_font_color(color):
    """设置字体颜色"""
    global current_paragraph
    for run in current_paragraph.runs:
        run.font.color.rgb = RGBColor.from_string(color2hex[color])


# 设置加粗字体
def set_font_bold():
    """设置字体加粗"""
    global current_paragraph
    for run in current_paragraph.runs:
        run.font.bold = True


# 设置斜体
def set_font_italic():
    """设置字体斜体"""
    global current_paragraph
    for run in current_paragraph.runs:
        run.font.italic = True


# 设置下划线
def set_font_underline():
    """设置字体下划线"""
    global current_paragraph
    for run in current_paragraph.runs:
        run.font.underline = True


# 设置字体样式
def set_font_style(font_name):
    """设置字体样式"""
    global current_paragraph
    for run in current_paragraph.runs:
        run.font.name = font_name


# 设置行间距
def set_line_space(line_space_level=0):
    """设置行间距"""
    global current_paragraph
    current_paragraph.paragraph_format.line_spacing = Pt(line_space_level)


# 左对齐
def text_align_left():
    """文本左对齐"""
    global current_paragraph
    current_paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT


# 居中对齐
def text_align_center():
    """文本居中对齐"""
    global current_paragraph
    current_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER


# 右对齐
def text_align_right():
    """文本右对齐"""
    global current_paragraph
    current_paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT


def add_run(text, font_size=12, font_color='000000', bold=False, italic=False, underline=False):
    global doc
    p = doc.paragraphs[-1]  # 添加到最新的段落
    run = p.add_run(text)
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor.from_string(font_color)
    run.font.bold = bold
    run.font.italic = italic
    run.font.underline = underline


def get_current_page_id():
    global doc
    # 在Word中，我们可以使用段落的数量来模拟页面编号
    # 因为Word文档不是分页存储的，所以这里的"页面编号"可能不完全准确
    # 但可以作为一个参考
    return len(doc.paragraphs)

def check_api_in_list_v2(api_call):
    api_names = [
    "set_word",
    "get_word",
    "create_docx",
    "save_word",
    "save_state",
    "load_state",
    "rename_docx_name",
    "move_docx_to_new_path",
    "get_current_page_id",
    "find_paragraphs_by_heading_and_content",
    "add_paragraph",
    "delete_paragraph",
    "add_text",
    "delete_text",
    "add_run",
    "add_line_break",
    "delete_line_break",
    "get_current_paragraph",
    "set_current_paragraph",
    "add_heading",
    "delete_heading",
    "add_table_of_contents",
    "delete_table_of_contents",
    "get_current_heading",
    "set_current_heading",
    "add_hyperlink",
    "delete_hyperlink",
    "add_header",
    "delete_header",
    "add_footer",
    "delete_footer",
    "add_page_numbers",
    "delete_page_numbers",
    "add_watermark",
    "delete_watermark",
    "add_table",
    "delete_table",
    "add_table_header",
    "set_table_title",
    "set_cell_text",
    "add_table_row",
    "delete_table_row",
    "add_table_column",
    "delete_table_column",
    "set_column_width",
    "set_row_height",
    "merge_cells",
    "set_cell_bg_color",
    "align_cell_text",
    "get_table_count",
    "get_current_table",
    "set_current_table",
    "add_list_item",
    "delete_list_item",
    "add_image",
    "delete_image",
    "set_image_size",
    "align_image_left",
    "align_image_right",
    "align_image_center",
    "add_caption",
    "replace_picture",
    "insert_picture",
    "get_current_picture",
    "set_current_picture",
    "insert_line_chart",
    "insert_bar_chart",
    "insert_pie_chart",
    "set_chart_title",
    "set_font_size",
    "set_font_color",
    "set_font_bold",
    "set_font_italic",
    "set_font_underline",
    "set_font_style",
    "set_line_space",
    "text_align_left",
    "text_align_center",
    "text_align_right",
    "check_api_in_list",
    "API_executor",
    "check_paragraph_before",
    "check_paragraph_after",
    "check_heading_level",
    "choose_paragraph",
    "choose_heading",
    "choose_table",
    "choose_table_cell",
    "choose_picture",
    "choose_object",
    "check"
]
    call_name = api_call.split('(')[0].strip()
    return call_name in api_names

def check_api_in_list(api_call, api_list):
    """检查API调用是否在指定的API列表中
    
    Args:
        api_call: API调用字符
        api_list: API名称列表
    Returns:
        bool: 是否在列表中
    """
    api_name = api_call.split('(')[0]
    return api_name in api_list

def API_executor(lines, test=False, args=None):
    """行API调用序列
    
    Args:
        lines: API调用序列（字符串列表）
        test: 是否为测试模式
        args: 参数配置
    Returns:
        error_info: 错误信息字符串
    """
    error_info = ""
    for line in lines:
        # 测试模式外的API过
        if not test:
            # 跳过位置相关API
            if check_api_in_list(line, ["set_left", "set_top", "set_right", "set_bottom"]):
                continue
            # API缺失模式下的过滤
            if args and args.api_lack:
                if not check_api_in_list(line, api_doc.original_apis):
                    continue
        
        try:
            # 长文档模式的特殊处理
            if args and args.dataset == 'long' and line == 'add_heading()':
                eval("add_paragraph('', style='Heading 1')")
                continue
            
            # API缺失模式
            if args and args.api_lack:
                if line == 'seek_assistance()':
                    continue
                elif not check_api_in_list(line, api_doc.original_apis):
                    # 添加标记段落
                    eval("add_paragraph('API不可用: @@@@@@@@@@')")
                else:
                    eval(line)
            
            # API更新模式
            elif args and args.api_update:
                if check_api_in_list(line, api_doc.update_apis):
                    # 添加标记段落
                    eval("add_paragraph('API已更新: @@@@@@@@@@')")
                else:
                    eval(line)
            
            # 正模式
            else:
                if not check_api_in_list_v2(line):
                    print(f"不存在的api调用: {line}")
                    eval("add_paragraph('API不可用: @@@@@@@@@@')")
                    continue
                eval(line)
                
        except Exception as e:
            error_msg = f"ERROR: {line}\n{str(e)}"
            print(error_msg)
            error_info += error_msg + "\n"
    
    return error_info

def get_current_paragraph():
    """获取当前操作的段落"""
    global current_paragraph
    return current_paragraph

def set_current_paragraph(paragraph):
    """设置当前操作的段落"""
    global current_paragraph
    current_paragraph = paragraph

def get_current_table():
    """获取当前操作的表格"""
    global current_table
    return current_table

def set_current_table(table):
    """设置当前操作的表格"""
    global current_table
    current_table = table

def get_current_picture():
    """获取当前操作的图片"""
    global current_picture
    return current_picture

def set_current_picture(picture):
    """设置当前操作的图片"""
    global current_picture
    current_picture = picture

def get_current_heading():
    """获取当前操作的标题"""
    global current_heading
    return current_heading

def set_current_heading(heading):
    """设置当前操作的标题"""
    global current_heading
    current_heading = heading

# 获取当前文档中已存在的表格数量
def get_table_count():
    """
    获取文档中表格的数量，返回已添加表格的数量。
    """
    global doc
    return len(doc.tables)

def add_table(rows, cols):
    """添加表格
    
    Args:
        rows: 行数
        cols: 列数
    Returns:
        新添加的表格对象
    """
    global doc, current_table
    try:
        current_table = doc.add_table(rows=rows, cols=cols)
        return current_table
    except Exception as e:
        logger.error(f"Error adding table: {e}")
        return None


# 设置表格标题
def set_table_title(title, font_size=12, bold=True, color="black", alignment=WD_PARAGRAPH_ALIGNMENT.CENTER):
    """
    设置表格标题并为表格添加编号。

    参数：
        doc (Document): docx 文档对象
        title (str): 表格标题文本
        font_size (int): 标题字体大小
        bold (bool): 是否加粗
        color (str): 标题字体颜色
        alignment (WD_PARAGRAPH_ALIGNMENT): 标题文本对齐方式
    """
    global doc
    table_count = get_table_count(doc) + 1  # 表格的编号是当前表格数 + 1
    table_title = f"Table {table_count}: {title}"

    # 在文档中添加标题段落，并设置居中对齐
    title_paragraph = doc.add_paragraph(table_title)
    title_paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER  # 设置标题居中

    # 设置标题样式
    set_font_size(font_size)
    set_font_color(color)
    if bold:
        set_font_bold()
    set_line_space(15)  # 设置标题的行间距


def delete_table(table_index):
    """
    删除指定索引的表格
    """
    global doc
    if 0 <= table_index < len(doc.tables):
        table = doc.tables[table_index]._element
        table.getparent().remove(table)


def add_table_header(table, headers, font_size=12, bold=True, alignment=WD_ALIGN_PARAGRAPH.CENTER):
    """
    添加或修改表格的表头行。

    参数：
        table (Table): 表格对象
        headers (list): 表头文本的列表
        font_size (int): 表头字体大小
        bold (bool): 是否加粗
        alignment (WD_ALIGN_PARAGRAPH): 表头文本的对齐方式
    """
    if all(cell.text == '' for cell in table.rows[0].cells):
        header_row = table.rows[0]
    else:
        header_row = table.add_row().cells
        table.rows.insert(0, table.rows.pop())

    for idx, header in enumerate(headers):
        cell = header_row.cells[idx]
        cell.text = header
        for paragraph in cell.paragraphs:
            paragraph.alignment = alignment
            run = paragraph.runs[0]
            run.font.size = Pt(font_size)
            run.font.bold = bold


def set_cell_text(row, col, text):
    """设置表格单元格文本
    
    Args:
        row: 行索引
        col: 列索引
        text: 单元格文本
    Returns:
        修改后的单元格对象
    """
    global current_table
    try:
        if current_table is None:
            logger.error("No table selected")
            return None
        cell = current_table.cell(row, col)
        cell.text = text
        return cell
    except Exception as e:
        logger.error(f"Error setting cell text: {e}")
        return None


def add_table_row(table_index):
    """
    向指定表格添加行
    """
    global doc
    if 0 <= table_index < len(doc.tables):
        doc.tables[table_index].add_row()

def add_table_column(table_index):
    """
    向指定表格添加列
    """
    global doc
    if 0 <= table_index < len(doc.tables):
        table = doc.tables[table_index]
        for row in table.rows:
            cell = row._tr.add_tc()
            table.cell(row._index, len(row.cells) - 1)._tc = cell
            cell = table.cell(row._index, len(row.cells) - 1)
            cell.width = Inches(1)  # 设置新列的宽度，可以根据需要调整


def delete_table_column(table_index, col_index):
    """
    删除表格中的指定列
    """
    global doc
    if 0 <= table_index < len(doc.tables):
        table = doc.tables[table_index]
        for row in table.rows:
            if 0 <= col_index < len(row.cells):
                tc = row.cells[col_index]._element
                tc.getparent().remove(tc)



def add_list_item(text, level=0, style='List Bullet'):
    """
    添加列表项

    参数:
    - level: 缩进级别
    - style: 列表样式('List Bullet'或'List Number')
    """
    global doc
    paragraph = doc.add_paragraph(style=style)
    paragraph.text = text
    paragraph.paragraph_format.left_indent = Pt(level * 18)


def delete_list_item(item_index):
    """
    删除指定的列表项及其所有子项
    """
    global doc
    count = 0
    to_delete = []
    current_level = None

    for i, paragraph in enumerate(doc.paragraphs):
        if paragraph.style.name.startswith('List'):
            if count == item_index:
                current_level = paragraph.paragraph_format.left_indent
                to_delete.append(i)
            elif current_level is not None:
                if paragraph.paragraph_format.left_indent > current_level:
                    to_delete.append(i)
                else:
                    break
            count += 1

    for i in reversed(to_delete):
        p = doc.paragraphs[i]._element
        p.getparent().remove(p)


def set_column_width(table, col_idx, width_in_inches):
    """
    设置表格某列的宽度

    参数：
        table (Table): 表格对象
        col_idx (int): 列索引
        width_in_inches (float): 列宽以英寸为单位）
    """
    for cell in table.columns[col_idx].cells:
        cell.width = Inches(width_in_inches)


def delete_table_row(table_index, row_index):
    """
    删除表格中的指定行
    """
    global doc
    if 0 <= table_index < len(doc.tables):
        table = doc.tables[table_index]
        if 0 <= row_index < len(table.rows):
            tr = table.rows[row_index]._element
            tr.getparent().remove(tr)


def add_table_column(table_index):
    """
    向指定表格添加列
    """
    global doc
    if 0 <= table_index < len(doc.tables):
        table = doc.tables[table_index]
        for row in table.rows:
            cell = row._tr.add_tc()
            table.cell(row._index, len(row.cells) - 1)._tc = cell
            cell = table.cell(row._index, len(row.cells) - 1)
            cell.width = Inches(1)  # 设置新列的宽度，可以根据需要调整


def set_row_height(table, row_idx, height_in_inches):
    """
    设置表格某行的高度。

    参数：
        table (Table): 表格对象
        row_idx (int): 行索引
        height_in_inches (float): 行高（以英寸为单位）
    """
    table.rows[row_idx].height = Inches(height_in_inches)


def merge_cells(table, start_row, start_col, end_row, end_col):
    """
    合并表格的单元格。

    参数：
        table (Table): 表格对象
        start_row (int): 起始行索引
        start_col (int): 起始列索引
        end_row (int): 结束行索引
        end_col (int): 结束列索引
    """
    cell = table.cell(start_row, start_col)
    other_cell = table.cell(end_row, end_col)
    cell.merge(other_cell)


def set_cell_bg_color(row, col, color):
    """设置表格单元格背景色
    
    Args:
        row: 行索引
        col: 列索引
        color: 颜色名称
    Returns:
        修改后的单元格对象
    """
    global current_table
    try:
        if current_table is None:
            logger.error("No table selected")
            return None
        cell = current_table.cell(row, col)
        # 设置背景色的实现...
        return cell
    except Exception as e:
        logger.error(f"Error setting cell background color: {e}")
        return None


def align_cell_text(table, row, col, alignment=WD_ALIGN_PARAGRAPH.CENTER):
    """
    设置单元格文本的对齐方式。

    参数：
        table (Table): 表格对象
        row (int): 行索引
        col (int): 列索引
        alignment (WD_ALIGN_PARAGRAPH): 对齐方式（如：WD_ALIGN_PARAGRAPH.CENTER、WD_ALIGN_PARAGRAPH.LEFT、WD_ALIGN_PARAGRAPH.RIGHT）
    """
    cell = table.cell(row, col)
    for paragraph in cell.paragraphs:
        paragraph.alignment = alignment


# picture


# 替换图片并调整大小的函数
def replace_picture(old_picture_path, new_picture_path, width_inch=None, height_inch=None):
    """
    替换文档中的旧图片为新图片，并调整图片的大小。
    :param doc: 文档对象
    :param old_picture_path: 旧图片的路径
    :param new_picture_path: 新图片的路径
    :param width_inch: 新图片的宽度，单位为英寸（可选）
    :param height_inch: 新图片的高度，单位为英寸（可选）
    """
    # 遍历文档中的段落，查找并删除旧图片
    global doc
    for paragraph in doc.paragraphs:
        for run in paragraph.runs:
            if run._r.pic:  # 检查当前run是否包含图片
                # 检查图片路径是否与旧图片路径匹配
                if old_picture_path in str(run._r.pic):
                    # 如果找到匹配的图片，清空该段落的内容（删除旧图片）
                    paragraph.clear()  # 清除该段落内容
                    break

    # 检查新图片是否存在
    if os.path.exists(new_picture_path):
        # 插入新图片
        picture = doc.add_picture(new_picture_path)

        # 如果指定了宽度，调整图片宽度
        if width_inch:
            picture.width = Inches(width_inch)

        # 如果指定了高度，调整图片高度
        if height_inch:
            picture.height = Inches(height_inch)
    else:
        print(f"错误: 找不到路为 {new_picture_path} 的文件。")


def add_image(paragraph_index, image_path, width=None, height=None, position=None):
    """
    在指定段落中添加图片
    
    参数:
    - paragraph_index: 段落索引
    - image_path: 图片文件路径
    - width: 图片宽度(单位为英寸)
    - height: 图片高度(单位为英寸)
    - position: 在段落中插入的位置(如果不指定,则在段落末尾添加)
    """
    global doc
    if 0 <= paragraph_index < len(doc.paragraphs):
        paragraph = doc.paragraphs[paragraph_index]
        
        # 创建新的 run 并移动到指定位置
        new_run = paragraph.add_run()
        if position is not None and 0 <= position < len(paragraph.runs):
            paragraph._p.insert(position, new_run._r)
            paragraph.runs.insert(position, new_run)
            paragraph.runs.pop()
        
        if width and height:
            new_run.add_picture(image_path, width=Inches(width), height=Inches(height))
        elif width:
            new_run.add_picture(image_path, width=Inches(width))
        elif height:
            new_run.add_picture(image_path, height=Inches(height))
        else:
            new_run.add_picture(image_path)



def delete_image(paragraph_index, image_index):
    """
    删除指定段落中的图片

    参数:
    - paragraph_index: 段落索引
    - image_index: 图片在段落中的索引
    """
    global doc
    if 0 <= paragraph_index < len(doc.paragraphs):
        paragraph = doc.paragraphs[paragraph_index]
        images = paragraph._element.findall('.//w:drawing', namespaces={
            'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'})
        if 0 <= image_index < len(images):
            image = images[image_index]
            parent = image.getparent()
            if parent is not None:
                parent.remove(image)
                # 如果 run 为空，则删除该 run
                if not parent.getchildren():
                    parent.getparent().remove(parent)


def set_image_size(image, width=None, height=None):
    """
    设置图片大小。

    :param image: 图片对象（在插入图片后返回的对象）
    :param width: 图片宽度（单位：英寸）
    :param height: 图片高度（单位：英寸）
    """
    if width:
        image.width = Inches(width)
    if height:
        image.height = Inches(height)


def align_image_left():
    """图片左对齐"""
    global current_paragraph
    try:
        if current_paragraph is None:
            logger.error("No paragraph selected")
            return None
        current_paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
        return current_paragraph
    except Exception as e:
        logger.error(f"Error aligning image left: {e}")
        return None


def align_image_right():
    """图片右对齐"""
    global current_paragraph
    try:
        if current_paragraph is None:
            logger.error("No paragraph selected")
            return None
        current_paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        return current_paragraph
    except Exception as e:
        logger.error(f"Error aligning image right: {e}")
        return None


def align_image_center():
    """图片居中对齐"""
    global current_paragraph
    try:
        if current_paragraph is None:
            logger.error("No paragraph selected")
            return None
        current_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        return current_paragraph
    except Exception as e:
        logger.error(f"Error aligning image center: {e}")
        return None


def add_caption(text, font_size=10):
    """
    为图片添加文字说明。


    :param text: 图片说明文字
    :param font_size: 字体大小（单位：Pt）
    """
    global doc
    caption_paragraph = doc.add_paragraph()
    run = caption_paragraph.add_run(text)
    run.font.size = Pt(font_size)


# chart

from docx.shared import Inches
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData

# 图表相关的全局变量
CHART_WIDTH = Inches(6)  # 图表默认宽度
CHART_HEIGHT = Inches(4)  # 图表默认高度
CHART_LEFT = Inches(1)   # 图表默认左边距
CHART_TOP = Inches(1)    # 图表默认上边距


def insert_picture(picture_name):
    """插入图片
    
    Args:
        picture_name: 图片名称(不含扩展名)
    """
    global doc, current_paragraph, current_picture
    try:
        # 创建新段落
        current_paragraph = doc.add_paragraph()
        current_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = current_paragraph.add_run()
        
        # 尝试插入图片
        try:
            current_picture = run.add_picture(f"{PIC_PATH}/{picture_name}.png")
        except:
            current_picture = run.add_picture(f"{PIC_PATH}/none.png")
            
        # 设置默认大小
        current_picture.width = Inches(6)
        current_picture.height = Inches(6)
        
        return current_picture
    except Exception as e:
        logger.error(f"Error inserting picture: {e}")
        return None

def insert_line_chart(x, y, title="Line Chart", xlabel="X Axis", ylabel="Y Axis"):
    """插入折线图
    
    Args:
        x: x轴标签列表
        y: y轴数据列表
        title: 图表标题
        xlabel: x轴标签
        ylabel: y轴标签
    """
    global doc, chart, current_paragraph
    try:
        # 创建新段落
        current_paragraph = doc.add_paragraph()
        current_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # 创建图表
        plt.figure(figsize=(8, 6))
        plt.plot(x, y, marker='o')
        plt.title(title)
        plt.xlabel(xlabel)
        plt.ylabel(ylabel)
        plt.grid(True)
        
        # 保存为临时文件
        temp_path = '_temp_line.png'
        plt.savefig(temp_path)
        plt.close()
        
        # 插入图片
        run = current_paragraph.add_run()
        chart = run.add_picture(temp_path)
        
        # 删除临时文件
        os.remove(temp_path)
        
        return chart
    except Exception as e:
        logger.error(f"Error inserting line chart: {e}")
        return None

def insert_bar_chart(x, y, title="Bar Chart", xlabel="X Axis", ylabel="Y Axis"):
    """插入柱状图"""
    global doc, chart, current_paragraph
    try:
        # 创建新段落
        current_paragraph = doc.add_paragraph()
        current_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # 创建图表
        plt.figure(figsize=(8, 6))
        plt.bar(x, y)
        plt.title(title, fontsize=12)
        plt.xlabel(xlabel, fontsize=10)
        plt.ylabel(ylabel, fontsize=10)
        plt.xticks(fontsize=9)
        plt.yticks(fontsize=9)
        plt.grid(True)
        
        # 保存为临时文件
        temp_path = '_temp_bar.png'
        plt.savefig(temp_path, dpi=300, bbox_inches='tight')
        plt.close()
        
        # 插入图片
        run = current_paragraph.add_run()
        chart = run.add_picture(temp_path)
        
        # 删除临时文件
        os.remove(temp_path)
        
        return chart
    except Exception as e:
        logger.error(f"Error inserting bar chart: {e}")
        return None

def insert_pie_chart(labels, sizes, title="Pie Chart"):
    """插入饼图
    
    Args:
        labels: 扇区标签列表
        sizes: 扇区大小列表
        title: 图表标题
    """
    global doc, chart, current_paragraph
    try:
        # 创建新段落
        current_paragraph = doc.add_paragraph()
        current_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # 创建图表
        plt.figure(figsize=(8, 8))
        plt.pie(sizes, labels=labels, autopct='%1.1f%%')
        plt.title(title)
        plt.axis('equal')
        
        # 保存为临时文件
        temp_path = '_temp_pie.png'
        plt.savefig(temp_path)
        plt.close()
        
        # 插入图片
        run = current_paragraph.add_run()
        chart = run.add_picture(temp_path)
        
        # 删除临时文件
        os.remove(temp_path)
        
        return chart
    except Exception as e:
        logger.error(f"Error inserting pie chart: {e}")
        return None

def set_chart_title(title):
    """设置图表标题
    
    Args:
        title: 标题文本
    """
    global chart
    try:
        if chart is None:
            logger.error("No chart selected")
            return None
        chart.chart_title.text_frame.text = title
        return chart
    except Exception as e:
        logger.error(f"Error setting chart title: {e}")
        return None
